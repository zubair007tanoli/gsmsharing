#!/usr/bin/env python3
"""
PDF to PowerPoint conversion: each PDF page becomes one slide (image).
Uses PyMuPDF (fitz) + python-pptx. pip install pymupdf python-pptx
"""
import sys
import os
import tempfile

def main(input_pdf, output_pptx):
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        print(f"ERROR: Install required packages: pip install pymupdf python-pptx. {e}", file=sys.stderr)
        sys.exit(1)

    try:
        if not os.path.exists(input_pdf):
            print(f"ERROR: Input file not found: {input_pdf}", file=sys.stderr)
            sys.exit(1)

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        doc = fitz.open(input_pdf)
        temp_dir = tempfile.mkdtemp()
        try:
            for i in range(len(doc)):
                page = doc[i]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                img_path = os.path.join(temp_dir, f"page_{i}.png")
                pix.save(img_path)
                blank = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            doc.close()
            out_dir = os.path.dirname(output_pptx)
            if out_dir and not os.path.isdir(out_dir):
                os.makedirs(out_dir, exist_ok=True)
            prs.save(output_pptx)
        finally:
            for f in os.listdir(temp_dir):
                try:
                    os.remove(os.path.join(temp_dir, f))
                except Exception:
                    pass
            try:
                os.rmdir(temp_dir)
            except Exception:
                pass

        if os.path.exists(output_pptx):
            print("SUCCESS")
            sys.exit(0)
        print("ERROR: Output file was not created.", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"ERROR: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("ERROR: Usage: pdf_to_pptx.py <input_pdf> <output_pptx>", file=sys.stderr)
        sys.exit(1)
    main(sys.argv[1], sys.argv[2])
